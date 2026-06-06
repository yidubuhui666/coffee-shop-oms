"""生成中文版精美答辩 PPT（咖啡店订单管理系统）。

设计要点：
  - 中文字体：微软雅黑（标题）/ 等线（正文）
  - 咖啡棕 + 米白配色，配合渐变与装饰几何元素
  - 卡片式布局、彩色色块、图标点缀
  - 16:9 宽屏，每页页脚分隔线 + 页码
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT  = r"C:\Users\许yq\coffee-shop-oms"
DOCS  = os.path.join(ROOT, "docs")
SHOTS = os.path.join(DOCS, "screenshots")
ER    = os.path.join(DOCS, "ER_Diagram.png")
OUT   = r"D:\cxdownload\092Xuyueqian-咖啡店订单管理系统.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

# ---------- 调色板 ----------
COFFEE_DARK   = RGBColor(0x2B, 0x1B, 0x14)   # 深咖
COFFEE       = RGBColor(0x6F, 0x4E, 0x37)    # 咖啡棕
COFFEE_LIGHT = RGBColor(0xA0, 0x66, 0x41)    # 浅咖
GOLD         = RGBColor(0xD7, 0xA8, 0x6E)    # 焦糖金
CREAM        = RGBColor(0xFA, 0xF4, 0xEC)    # 奶白
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT         = RGBColor(0x2A, 0x21, 0x1B)
TEXT_SOFT    = RGBColor(0x6B, 0x55, 0x47)
GRAY         = RGBColor(0x8B, 0x75, 0x69)
GREEN        = RGBColor(0x4F, 0x8B, 0x52)
RED          = RGBColor(0xB0, 0x48, 0x35)
BLUE         = RGBColor(0x3C, 0x6E, 0x9F)

FONT_H = "微软雅黑"
FONT_B = "等线"
FONT_M = "Consolas"  # 等宽

blank = prs.slide_layouts[6]


# ---------- 工具 ----------
def add_rect(slide, left, top, width, height, fill=None, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if not line:
        shp.line.fill.background()
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.shadow.inherit = False
    return shp


def add_round(slide, left, top, width, height, fill, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.adjustments[0] = radius
    shp.line.fill.background()
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font=FONT_B, anchor=None, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # 中文字体回退
    try:
        rPr = run._r.get_or_add_rPr()
        from pptx.oxml.ns import qn
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            from lxml import etree
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font)
    except Exception:
        pass
    return box


def add_bullets(slide, items, left, top, width, height,
                size=16, color=TEXT, line=1.4, bullet_color=None,
                font=FONT_B, dot="●"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line
        # 项目符号
        r1 = p.add_run(); r1.text = dot + "  "
        r1.font.size = Pt(size - 2)
        r1.font.color.rgb = bullet_color or GOLD
        r1.font.name = font
        # 内容
        r2 = p.add_run(); r2.text = it
        r2.font.size = Pt(size); r2.font.color.rgb = color
        r2.font.name = font


def footer(slide, page, total):
    # 分隔线
    line = add_rect(slide, Inches(0.5), Inches(7.05),
                    Inches(12.3), Emu(9525), fill=GOLD)
    add_text(slide, "咖啡店订单管理系统  ·  数据库原理与应用  ·  西南大学",
             Inches(0.5), Inches(7.12), Inches(10), Inches(0.3),
             size=10, color=GRAY)
    add_text(slide, f"{page} / {total}", Inches(11.5), Inches(7.12),
             Inches(1.3), Inches(0.3),
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def title_bar(slide, no, title, subtitle=None):
    # 左侧色块
    add_round(slide, Inches(0.5), Inches(0.5), Inches(0.18), Inches(0.85),
              GOLD, radius=0.3)
    # 编号
    add_text(slide, no, Inches(0.78), Inches(0.5), Inches(1.4), Inches(0.6),
             size=14, color=GOLD, bold=True, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    # 主标题
    add_text(slide, title, Inches(0.78), Inches(0.85), Inches(11), Inches(0.6),
             size=26, bold=True, color=COFFEE_DARK, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.78), Inches(1.45),
                 Inches(11), Inches(0.4),
                 size=12, color=GRAY, font=FONT_B)


def card(slide, left, top, width, height, fill=WHITE):
    """阴影白卡片（用淡灰矩形模拟阴影）。"""
    # 阴影
    sh = add_round(slide, left + Emu(38100), top + Emu(38100),
                   width, height, RGBColor(0xE6, 0xDE, 0xD2), radius=0.04)
    c = add_round(slide, left, top, width, height, fill, radius=0.04)
    return c


# ====================================================================
# 1. 封面
# ====================================================================
def slide_cover():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=COFFEE_DARK)

    # 左侧大色块 + 装饰
    add_rect(s, 0, 0, Inches(6.5), SH, fill=COFFEE)
    # 装饰金线
    add_rect(s, Inches(0.7), Inches(2.6), Inches(0.8), Emu(38100), fill=GOLD)

    # 大标题
    add_text(s, "咖啡店", Inches(0.7), Inches(2.85), Inches(6),
             Inches(1.4), size=72, bold=True, color=CREAM, font=FONT_H)
    add_text(s, "订单管理系统", Inches(0.7), Inches(4.05), Inches(6),
             Inches(1.0), size=44, bold=True, color=GOLD, font=FONT_H)
    add_text(s, "Coffee Shop Order Management System",
             Inches(0.7), Inches(5.1), Inches(6), Inches(0.5),
             size=14, color=CREAM, font=FONT_B)

    # 副标
    add_text(s, "《数据库原理与应用》课程设计",
             Inches(0.7), Inches(0.7), Inches(6), Inches(0.5),
             size=15, color=GOLD, font=FONT_H)

    # 右侧信息卡
    add_round(s, Inches(7.1), Inches(2.3), Inches(5.6), Inches(3.0),
              CREAM, radius=0.05)
    add_text(s, "答辩信息", Inches(7.4), Inches(2.5),
             Inches(5), Inches(0.4),
             size=14, color=GOLD, bold=True, font=FONT_H)
    add_text(s, "小组成员", Inches(7.4), Inches(2.95),
             Inches(5), Inches(0.4),
             size=12, color=GRAY, font=FONT_B)
    add_text(s, "114 杨帆（组长）   ·   092 许跃骞",
             Inches(7.4), Inches(3.25), Inches(5), Inches(0.4),
             size=18, bold=True, color=COFFEE_DARK, font=FONT_H)

    add_text(s, "专业 / 年级", Inches(7.4), Inches(3.85),
             Inches(5), Inches(0.4),
             size=12, color=GRAY, font=FONT_B)
    add_text(s, "计算机科学与技术 · 2024级",
             Inches(7.4), Inches(4.15), Inches(5), Inches(0.4),
             size=16, color=COFFEE_DARK, font=FONT_H)

    add_text(s, "学院", Inches(7.4), Inches(4.7),
             Inches(5), Inches(0.4),
             size=12, color=GRAY, font=FONT_B)
    add_text(s, "西南大学  ·  计算机与信息科学学院",
             Inches(7.4), Inches(4.95), Inches(5), Inches(0.4),
             size=14, color=COFFEE_DARK, font=FONT_H)

    # 底部装饰
    add_rect(s, 0, Inches(7.2), SW, Inches(0.3), fill=GOLD)
    add_text(s, "☕  2026 年 6 月", 0, Inches(7.2),
             SW, Inches(0.3), size=11, color=COFFEE_DARK,
             align=PP_ALIGN.CENTER, font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# 2. 目录
# ====================================================================
def slide_outline():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    # 左大字
    add_text(s, "CONTENTS", Inches(0.6), Inches(1.0), Inches(5), Inches(0.6),
             size=14, color=GOLD, bold=True, font=FONT_H)
    add_text(s, "目  录", Inches(0.6), Inches(1.4), Inches(5), Inches(1.4),
             size=60, bold=True, color=COFFEE_DARK, font=FONT_H)
    # 金色装饰条
    add_rect(s, Inches(0.6), Inches(2.8), Inches(0.8), Emu(57150), fill=GOLD)

    # 右侧目录条目（卡片）
    items = [
        ("01", "分工说明",       "Division of Labor"),
        ("02", "项目概述",       "Project Overview"),
        ("03", "数据库设计",     "Database Design & SQL"),
        ("04", "使用说明",       "Brief Instructions"),
        ("05", "界面展示",       "Interface Screenshots"),
        ("06", "总结与展望",     "Summary & Future Work"),
        ("07", "参考文献",       "References"),
    ]
    base_top = Inches(0.9)
    row_h = Inches(0.78)
    for i, (no, zh, en) in enumerate(items):
        top = base_top + row_h * i
        # 编号圆角块
        add_round(s, Inches(7.0), top, Inches(0.85), Inches(0.65),
                  COFFEE if i % 2 == 0 else GOLD, radius=0.25)
        add_text(s, no, Inches(7.0), top, Inches(0.85), Inches(0.65),
                 size=18, bold=True, color=CREAM, align=PP_ALIGN.CENTER,
                 font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
        # 中文标题
        add_text(s, zh, Inches(8.05), top + Emu(50000),
                 Inches(4), Inches(0.4),
                 size=20, bold=True, color=COFFEE_DARK, font=FONT_H)
        # 英文小字
        add_text(s, en, Inches(8.05), top + Inches(0.36),
                 Inches(4), Inches(0.3),
                 size=11, color=GRAY, font=FONT_B)
    footer(s, 2, TOTAL)


# ====================================================================
# 3. 分工
# ====================================================================
def slide_division():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "01", "分工说明",
              "两人小组 · 任务粒度协作分配 · Git 版本控制")

    # 两个成员卡片
    def member_card(left, name, sid, role, color, tasks):
        card(s, left, Inches(2.0), Inches(5.8), Inches(4.5))
        # 顶部色条
        add_rect(s, left, Inches(2.0), Inches(5.8),
                 Inches(0.12), fill=color)
        # 头像圆
        add_round(s, left + Inches(0.35), Inches(2.3),
                  Inches(0.95), Inches(0.95), color, radius=0.5)
        add_text(s, name[0], left + Inches(0.35), Inches(2.3),
                 Inches(0.95), Inches(0.95),
                 size=30, bold=True, color=CREAM,
                 align=PP_ALIGN.CENTER, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 姓名
        add_text(s, name, left + Inches(1.5), Inches(2.35),
                 Inches(4), Inches(0.5),
                 size=22, bold=True, color=COFFEE_DARK, font=FONT_H)
        add_text(s, f"学号 {sid}  ·  {role}",
                 left + Inches(1.5), Inches(2.85),
                 Inches(4), Inches(0.4),
                 size=12, color=GRAY, font=FONT_B)
        # 任务列表
        add_bullets(s, tasks,
                    left + Inches(0.4), Inches(3.55),
                    Inches(5.2), Inches(2.9),
                    size=14, line=1.6, bullet_color=color)

    member_card(Inches(0.5), "杨  帆", "114", "组长 · 全栈参与（偏后端 / 数据库）", COFFEE, [
        "需求分析 + ER 建模主导",
        "Schema · 触发器 · 视图 · 存储过程",
        "后端路由：订单 / 库存 / 统计模块",
        "前端模板：仪表盘 / 订单管理",
        "演示数据 bootstrap 脚本",
    ])
    member_card(Inches(7.05), "许跃骞", "092", "组员 · 全栈参与（偏前端 / 应用）", GOLD, [
        "后端路由：菜单 / 会员 / 促销模块",
        "前端模板与 CSS 主样式",
        "AI 生成商品图片素材",
        "一键启动脚本（.bat / .command）",
        "答辩报告 + PPT + 数据库迁移脚本",
    ])

    # 协作方式
    add_text(s, "🔗  协作方式：两人各自维护独立 MySQL 数据库  ·  通过 GitHub 同步代码  ·  数据库改动配套幂等种子脚本（git pull 即同步）",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.35),
             size=11, color=COFFEE, font=FONT_H, align=PP_ALIGN.CENTER)
    footer(s, 3, TOTAL)


# ====================================================================
# 4. 项目概述
# ====================================================================
def slide_overview():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "02", "项目概述",
              "数字化一家独立咖啡店的日常运营 · 取代纸质点单与表格记账")

    # 三栏卡片：背景 / 任务 / 难点
    def info_card(left, top, w, h, icon, color, title, lines):
        card(s, left, top, w, h)
        # 顶部色块带图标
        add_rect(s, left, top, w, Inches(0.6), fill=color)
        add_text(s, icon, left + Inches(0.2), top,
                 Inches(0.6), Inches(0.6),
                 size=22, color=CREAM, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, title, left + Inches(0.85), top,
                 w - Inches(1), Inches(0.6),
                 size=18, bold=True, color=CREAM, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 内容
        add_bullets(s, lines, left + Inches(0.3), top + Inches(0.85),
                    w - Inches(0.6), h - Inches(1),
                    size=13, line=1.5, bullet_color=color)

    info_card(Inches(0.5), Inches(2.0), Inches(4.1), Inches(4.7),
              "📋", COFFEE, "项目背景", [
        "灵感来自学校旁边的独立咖啡店",
        "店主原本用纸质点单 + Excel 记账",
        "易出错、无任何数据分析能力",
        "目标：用数据库 + Web 端替代",
    ])
    info_card(Inches(4.7), Inches(2.0), Inches(4.1), Inches(4.7),
              "🎯", GOLD, "主要任务", [
        "完整描述咖啡店业务的关系模式",
        "在数据库中实现完整性、触发器等",
        "三种角色（管理员/咖啡师/收银员）",
        "对所有实体提供完整 CRUD",
        "丰富的多表 JOIN 报表",
    ])
    info_card(Inches(8.9), Inches(2.0), Inches(3.9), Inches(4.7),
              "🧩", COFFEE_LIGHT, "关键难点", [
        "1对多 / 多对多关系建模",
        "订单总额随商品项联动一致",
        "积分自动入账 + 兑换抵现",
        "原料库存随下单消耗 + 预警",
        "会员折扣 ⊕ 促销码取更优",
    ])
    footer(s, 4, TOTAL)


# ====================================================================
# 5. ER 图
# ====================================================================
def slide_er():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·1", "概念设计  —  ER 图",
              "9 个实体 + 1 个联系实体 ProductIngredient · 三类关系并存")
    card(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(4.95), WHITE)
    if os.path.exists(ER):
        s.shapes.add_picture(ER, Inches(1.5), Inches(2.05),
                             height=Inches(4.75))

    # 右下角图例
    add_text(s, "实体类型",
             Inches(9.5), Inches(2.1), Inches(3.0), Inches(0.4),
             size=12, bold=True, color=GOLD, font=FONT_H)
    legend = [
        ("customers / staff",       "强实体"),
        ("orders / order_items",    "事件实体"),
        ("products / inventory",    "资源实体"),
        ("product_ingredient",      "联系实体（M:N）"),
    ]
    for i, (name, kind) in enumerate(legend):
        add_text(s, "● " + name,
                 Inches(9.5), Inches(2.5) + Inches(0.45) * i,
                 Inches(3.2), Inches(0.3), size=11, color=COFFEE_DARK,
                 font=FONT_M)
        add_text(s, "  " + kind,
                 Inches(9.5), Inches(2.7) + Inches(0.45) * i,
                 Inches(3.2), Inches(0.3), size=10, color=GRAY,
                 font=FONT_B)
    footer(s, 6, TOTAL)


# ====================================================================
# 6. 逻辑设计 - 10个关系
# ====================================================================
def slide_logical():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·2", "逻辑设计  —  十张关系表",
              "已规范化至 3NF / BCNF · PK 主键 · FK 外键")

    rels = [
        ("customers",          "会员信息（积分、等级）"),
        ("staff",              "员工（角色、密码）"),
        ("categories",         "商品分类"),
        ("products",           "商品菜单"),
        ("promotions",         "促销码"),
        ("orders",             "订单主表"),
        ("order_items",        "订单明细"),
        ("inventory",          "原料库存"),
        ("product_ingredient", "配方（M:N 联系表）"),
        ("member_discount",    "会员等级折扣"),
    ]

    cols = 2
    cw, ch = Inches(6.0), Inches(0.85)
    base_left, base_top = Inches(0.6), Inches(2.0)
    for i, (name, desc) in enumerate(rels):
        r, c = i // cols, i % cols
        left = base_left + (cw + Inches(0.15)) * c
        top  = base_top + (ch + Inches(0.10)) * r
        card(s, left, top, cw, ch, WHITE)
        # 序号
        add_round(s, left + Inches(0.18), top + Inches(0.17),
                  Inches(0.5), Inches(0.5), COFFEE, radius=0.25)
        add_text(s, str(i + 1),
                 left + Inches(0.18), top + Inches(0.17),
                 Inches(0.5), Inches(0.5),
                 size=14, bold=True, color=CREAM,
                 align=PP_ALIGN.CENTER, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 表名
        add_text(s, name, left + Inches(0.85), top + Inches(0.1),
                 Inches(5), Inches(0.4),
                 size=15, bold=True, color=COFFEE_DARK, font=FONT_M)
        # 描述
        add_text(s, desc, left + Inches(0.85), top + Inches(0.42),
                 Inches(5), Inches(0.4),
                 size=12, color=GRAY, font=FONT_B)
    footer(s, 7, TOTAL)


# ====================================================================
# 7. DDL 摘录
# ====================================================================
def slide_ddl():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·3", "物理设计  —  DDL 摘录",
              "完整 DDL 见 sql/schema.sql · 此处展示 orders 表")

    code = (
        "CREATE TABLE orders (\n"
        "    order_id        INT AUTO_INCREMENT PRIMARY KEY,\n"
        "    customer_id     INT,\n"
        "    staff_id        INT NOT NULL,\n"
        "    promo_id        INT,\n"
        "    order_time      DATETIME DEFAULT CURRENT_TIMESTAMP,\n"
        "    total_amount    DECIMAL(10,2) NOT NULL DEFAULT 0,\n"
        "    status          ENUM('PENDING','PAID','PREPARING','READY',\n"
        "                         'COMPLETED','CANCELLED'),\n"
        "    payment_method  ENUM('CASH','CARD','WECHAT','ALIPAY'),\n"
        "    points_used     INT DEFAULT 0,\n"
        "    FOREIGN KEY (customer_id) REFERENCES customers(customer_id)\n"
        "                  ON DELETE SET NULL,\n"
        "    FOREIGN KEY (staff_id)    REFERENCES staff(staff_id),\n"
        "    FOREIGN KEY (promo_id)    REFERENCES promotions(promo_id),\n"
        "    INDEX idx_time(order_time),\n"
        "    INDEX idx_status(status)\n"
        ") ENGINE=InnoDB;"
    )
    card(s, Inches(0.5), Inches(2.0), Inches(8.5), Inches(4.9),
         RGBColor(0x2B, 0x1B, 0x14))
    box = s.shapes.add_textbox(Inches(0.75), Inches(2.15),
                               Inches(8.1), Inches(4.6))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        run = p.add_run(); run.text = line
        run.font.name = FONT_M
        run.font.size = Pt(11)
        # 关键字高亮
        kw = ("CREATE", "TABLE", "PRIMARY", "KEY", "NOT", "NULL",
              "DEFAULT", "FOREIGN", "REFERENCES", "ON", "DELETE", "SET",
              "INDEX", "INT", "ENGINE", "AUTO_INCREMENT", "DATETIME",
              "CURRENT_TIMESTAMP", "DECIMAL", "ENUM")
        if any(k in line for k in kw):
            run.font.color.rgb = GOLD
        else:
            run.font.color.rgb = CREAM

    # 右侧亮点
    card(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.9), WHITE)
    add_text(s, "设计亮点", Inches(9.4), Inches(2.15),
             Inches(3.4), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)
    add_bullets(s, [
        "ENUM 限定订单状态合法取值",
        "外键 ON DELETE SET NULL\n保留历史订单完整性",
        "order_time / status 加索引\n加速仪表盘与筛选",
        "DECIMAL(10,2) 保证金额精度",
        "InnoDB 引擎，支持事务",
    ], Inches(9.4), Inches(2.6), Inches(3.4), Inches(4.2),
       size=11, line=1.45)
    footer(s, 8, TOTAL)


# ====================================================================
# 8. 触发器 + 存储过程
# ====================================================================
def slide_triggers():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·4", "数据库编程  —  触发器 + 存储过程 + 视图",
              "把业务规则放进数据库 · 与客户端无关地保证不变量")

    # 4 个小卡片
    items = [
        ("①", "trg_item_after_insert",   "插入订单项 → 自动重算 total_amount",        COFFEE),
        ("②", "trg_order_after_update",  "状态变 PAID → 自动入账积分（1 RMB ⇒ 1 分）", GOLD),
        ("③", "trg_item_deduct_stock",   "下单按配方扣减原料库存",                     COFFEE_LIGHT),
        ("④", "sp_place_order",          "一次调用完成单品下单（IN 4 / OUT 1）",       BLUE),
    ]
    cw, ch = Inches(6.0), Inches(1.0)
    for i, (no, name, desc, color) in enumerate(items):
        left = Inches(0.5) + (cw + Inches(0.15)) * (i % 2)
        top  = Inches(2.0) + (ch + Inches(0.15)) * (i // 2)
        card(s, left, top, cw, ch, WHITE)
        add_rect(s, left, top, Inches(0.15), ch, fill=color)
        add_text(s, no, left + Inches(0.3), top, Inches(0.6), ch,
                 size=26, bold=True, color=color,
                 font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, left + Inches(1.0), top + Inches(0.1),
                 Inches(5), Inches(0.4),
                 size=14, bold=True, color=COFFEE_DARK, font=FONT_M)
        add_text(s, desc, left + Inches(1.0), top + Inches(0.5),
                 Inches(5), Inches(0.4),
                 size=12, color=TEXT_SOFT, font=FONT_B)

    # 触发器代码示例
    card(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4),
         RGBColor(0x2B, 0x1B, 0x14))
    code = (
        "-- 触发器示例 ②：订单变更为 PAID 时自动入账会员积分\n"
        "CREATE TRIGGER trg_order_after_update AFTER UPDATE ON orders\n"
        "FOR EACH ROW BEGIN\n"
        "    IF NEW.status = 'PAID' AND OLD.status <> 'PAID'\n"
        "       AND NEW.customer_id IS NOT NULL THEN\n"
        "        UPDATE customers SET points = points + FLOOR(NEW.total_amount)\n"
        "        WHERE  customer_id = NEW.customer_id;\n"
        "    END IF;\n"
        "END;"
    )
    box = s.shapes.add_textbox(Inches(0.75), Inches(4.6),
                               Inches(11.8), Inches(2.2))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        r = p.add_run(); r.text = line
        r.font.name = FONT_M; r.font.size = Pt(12)
        r.font.color.rgb = (GOLD if i == 0 else CREAM)
    footer(s, 12, TOTAL)


# ====================================================================
# 9. 代表性 SQL
# ====================================================================
def slide_queries():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·5", "代表性 SQL 查询",
              "全部出现在应用的实际页面里 · 充分体现多表关系")

    blocks = [
        ("当日营收（聚合 + 过滤）",
         "SELECT SUM(total_amount) FROM orders\n"
         "WHERE  DATE(order_time) = CURDATE()\n"
         "  AND  status <> 'CANCELLED';"),
        ("商品销量 TOP-5（2 表 JOIN）",
         "SELECT p.name, SUM(oi.quantity) qty\n"
         "FROM   order_items oi JOIN products p USING(product_id)\n"
         "GROUP  BY p.product_id ORDER BY qty DESC LIMIT 5;"),
        ("会员消费排行（3 表 JOIN + 聚合）",
         "SELECT c.name, c.member_level,\n"
         "       SUM(oi.quantity * oi.unit_price) spent\n"
         "FROM   customers c\n"
         "JOIN   orders o      ON o.customer_id = c.customer_id\n"
         "JOIN   order_items oi ON oi.order_id   = o.order_id\n"
         "GROUP  BY c.customer_id ORDER BY spent DESC LIMIT 10;"),
        ("原料消耗预测（4 表 LEFT JOIN）",
         "SELECT inv.name, SUM(pi.consume_qty * oi.quantity) used\n"
         "FROM   inventory inv\n"
         "LEFT JOIN product_ingredient pi USING(ingredient_id)\n"
         "LEFT JOIN order_items oi USING(product_id)\n"
         "LEFT JOIN orders o USING(order_id)\n"
         "WHERE  o.order_time >= DATE_SUB(CURDATE(), INTERVAL 7 DAY);"),
    ]

    cw, ch = Inches(6.05), Inches(2.4)
    for i, (title, code) in enumerate(blocks):
        left = Inches(0.5) + (cw + Inches(0.15)) * (i % 2)
        top  = Inches(2.0) + (ch + Inches(0.13)) * (i // 2)
        card(s, left, top, cw, ch, WHITE)
        # 标题条
        add_rect(s, left, top, cw, Inches(0.42), fill=COFFEE)
        add_text(s, title, left + Inches(0.2), top,
                 cw - Inches(0.3), Inches(0.42),
                 size=12, bold=True, color=CREAM, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 代码
        box = s.shapes.add_textbox(left + Inches(0.15),
                                   top + Inches(0.5),
                                   cw - Inches(0.3),
                                   ch - Inches(0.55))
        tf = box.text_frame; tf.word_wrap = True
        for j, line in enumerate(code.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.15
            r = p.add_run(); r.text = line
            r.font.name = FONT_M; r.font.size = Pt(10)
            r.font.color.rgb = COFFEE_DARK
    footer(s, 13, TOTAL)


# ====================================================================
# 10. 使用说明 + 账号
# ====================================================================
def slide_usage():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "04", "使用说明",
              "一键启动 · 跨平台双脚本 · 4 个预置演示账号")

    # 左：环境 / 启动
    card(s, Inches(0.5), Inches(2.0), Inches(6.3), Inches(4.8), WHITE)
    add_text(s, "🚀 启动方式", Inches(0.8), Inches(2.15),
             Inches(5.5), Inches(0.5),
             size=18, bold=True, color=GOLD, font=FONT_H)
    add_text(s, "环境要求",
             Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.4),
             size=13, bold=True, color=COFFEE_DARK, font=FONT_H)
    add_bullets(s, [
        "Python 3.10+    ·   MySQL 8.0+",
        "依赖：Flask · Flask-SQLAlchemy · PyMySQL",
    ], Inches(0.8), Inches(3.05), Inches(5.5), Inches(1.0),
       size=12, line=1.45)

    add_text(s, "一键启动",
             Inches(0.8), Inches(4.1), Inches(5.5), Inches(0.4),
             size=13, bold=True, color=COFFEE_DARK, font=FONT_H)
    add_bullets(s, [
        "Windows ：双击  启动.bat",
        "macOS  ：双击  启动咖啡店系统.command",
        "自动建虚拟环境 · 装依赖 · 起服务",
        "浏览器自动打开 http://127.0.0.1:5050",
    ], Inches(0.8), Inches(4.45), Inches(5.5), Inches(2.3),
       size=12, line=1.5)

    # 右：账号表
    card(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.8), WHITE)
    add_text(s, "👥 演示账号", Inches(7.3), Inches(2.15),
             Inches(5), Inches(0.5),
             size=18, bold=True, color=GOLD, font=FONT_H)
    rows = [
        ("角色",    "用户名",   "密码",         "姓名"),
        ("管理员",  "admin",   "admin123",    "杨帆"),
        ("管理员",  "xyq",     "070203",      "许跃骞"),
        ("咖啡师",  "barista", "barista123",  "李明"),
        ("收银员",  "cashier", "cashier123",  "王磊"),
    ]
    table = s.shapes.add_table(rows=len(rows), cols=4,
                               left=Inches(7.25), top=Inches(2.8),
                               width=Inches(5.3),
                               height=Inches(3.6)).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j); cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT_H if i == 0 else FONT_B
                    r.font.size = Pt(13)
                    r.font.bold = (i == 0)
                    r.font.color.rgb = CREAM if i == 0 else COFFEE_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (COFFEE if i == 0 else
                (RGBColor(0xFA, 0xF4, 0xEC) if i % 2 else WHITE))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer(s, 15, TOTAL)


# ====================================================================
# 11-15. 截图
# ====================================================================
SHOTS_META = [
    ("01_dashboard.png", "仪表盘",       "KPI 卡片 · 当日营收 · 商品销量 TOP-5"),
    ("02_menu.png",      "菜单页",       "分类化商品卡片 · 关键词与分类筛选"),
    ("03_new_order.png", "下单页面",     "多行购物车 · 规格 · 支付方式 · 折扣 / 积分抵现"),
    ("04_orders.png",    "订单列表",     "今日/历史标签 · 状态筛选 · 分页 · 状态机流转"),
    ("05_stats.png",     "统计分析",     "营收趋势 · 销量排行 · 8 个多表 JOIN 报表"),
]

def slide_screenshot(idx, fname, name, desc, page):
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, f"05·{idx}", f"界面展示  —  {name}", desc)
    # 白卡放截图
    card(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.9), WHITE)
    path = os.path.join(SHOTS, fname)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(1.0), Inches(2.15),
                             width=Inches(11.3))
    footer(s, page, TOTAL)


# ====================================================================
# 16. 工程改进
# ====================================================================
def slide_engineering():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "06·2", "工程改进",
              "迭代过程中真实落地的代码质量提升")

    items = [
        ("🧩",  "Blueprint 模块化重构",
                "1000 行单文件 app.py  →  11 个 Blueprint 模块（应用工厂模式）",
                COFFEE),
        ("⚡",  "消除 N+1 查询",
                "joinedload / selectinload  · 单页 SQL 数：~70 → ~5",
                GOLD),
        ("📄",  "服务端分页",
                "订单列表 / 会员列表，避免一次加载全表",
                COFFEE_LIGHT),
        ("🛡️", "友好错误页 + 冒烟测试",
                "404 / 500 自定义页面 · 4 项冒烟测试覆盖核心路径",
                BLUE),
        ("🔄",  "数据库改动配套幂等脚本",
                "每次改 DB 都附 seed_*.py / migration_*.py，搭档拉取即同步",
                GREEN),
        ("🎨",  "AI 生成商品图 + 演示数据",
                "29 张商品图 · 20 位会员 · 150 条历史订单",
                RED),
    ]
    cw, ch = Inches(6.05), Inches(1.4)
    for i, (icon, title, desc, color) in enumerate(items):
        left = Inches(0.5) + (cw + Inches(0.15)) * (i % 2)
        top  = Inches(2.0) + (ch + Inches(0.13)) * (i // 2)
        card(s, left, top, cw, ch, WHITE)
        add_rect(s, left, top, Inches(0.15), ch, fill=color)
        add_text(s, icon, left + Inches(0.25), top,
                 Inches(0.9), ch,
                 size=28, color=color,
                 font=FONT_H, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        add_text(s, title, left + Inches(1.2), top + Inches(0.2),
                 cw - Inches(1.3), Inches(0.4),
                 size=15, bold=True, color=COFFEE_DARK, font=FONT_H)
        add_text(s, desc, left + Inches(1.2), top + Inches(0.65),
                 cw - Inches(1.3), Inches(0.7),
                 size=11, color=TEXT_SOFT, font=FONT_B,
                 line_spacing=1.3)
    footer(s, 24, TOTAL)


# ====================================================================
# 17. 现存问题 / 展望
# ====================================================================
def slide_future():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "06·3", "现存问题与未来改进",
              "已识别但暂未优化的点 · 后续迭代方向")

    items = [
        ("🔐", "密码哈希存储",        "目前明文存储 · 应改为 bcrypt 或 werkzeug.security 哈希"),
        ("🚀", "生产级部署",          "单 Flask 进程 · 应改为 Gunicorn + Nginx 反向代理"),
        ("📊", "可视化升级",          "纯模板渲染 · 引入 ECharts 让统计页更直观"),
        ("📝", "操作日志表",          "增加审计表（谁、何时、改了什么）"),
        ("📱", "移动端适配",          "导航栏在窄屏挤压 · 加汉堡菜单或换行"),
        ("🔍", "数据库索引补全",       "0 个自定义索引 · 给热点外键列加索引"),
    ]
    for i, (icon, title, desc) in enumerate(items):
        top = Inches(2.0) + Inches(0.78) * i
        card(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
        add_text(s, icon, Inches(0.65), top,
                 Inches(0.7), Inches(0.7),
                 size=22, color=GOLD, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(1.4), top + Inches(0.1),
                 Inches(3.2), Inches(0.5),
                 size=15, bold=True, color=COFFEE_DARK, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(4.6), top + Inches(0.1),
                 Inches(8), Inches(0.5),
                 size=12, color=TEXT_SOFT, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 25, TOTAL)


# ====================================================================
# 18. Thank you
# ====================================================================
def slide_thanks():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=COFFEE_DARK)
    # 大装饰圆
    add_round(s, Inches(-2), Inches(-2), Inches(5), Inches(5),
              COFFEE, radius=0.5)
    add_round(s, Inches(10), Inches(4), Inches(6), Inches(6),
              COFFEE_LIGHT, radius=0.5)
    # 金线
    add_rect(s, Inches(5.5), Inches(3.0), Inches(2.3),
             Emu(57150), fill=GOLD)
    # 大字
    add_text(s, "感  谢  聆  听", 0, Inches(2.2),
             SW, Inches(1.2),
             size=58, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, font=FONT_H)
    add_text(s, "THANK YOU", 0, Inches(3.4),
             SW, Inches(0.6),
             size=20, color=GOLD,
             align=PP_ALIGN.CENTER, font=FONT_H)
    add_text(s, "敬请提问", 0, Inches(4.6),
             SW, Inches(0.6),
             size=24, color=CREAM,
             align=PP_ALIGN.CENTER, font=FONT_H)
    add_text(s, "Q  &  A", 0, Inches(5.2), SW, Inches(0.5),
             size=18, color=GOLD,
             align=PP_ALIGN.CENTER, font=FONT_H)
    # 底栏
    add_rect(s, 0, Inches(6.8), SW, Inches(0.7), fill=GOLD)
    add_text(s, "☕  咖啡店订单管理系统  ·  114 杨帆 & 092 许跃骞  ·  西南大学",
             0, Inches(6.8), SW, Inches(0.7),
             size=14, color=COFFEE_DARK, bold=True, font=FONT_H,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# 新增：技术栈 + 系统架构
# ====================================================================
def slide_techstack():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "02·补", "技术栈与系统架构",
              "三层 B/S 架构 · Python 全栈 · 浏览器零安装即用")

    # 左：技术栈表格
    card(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.8), WHITE)
    add_text(s, "🛠️ 技术栈", Inches(0.8), Inches(2.15),
             Inches(5), Inches(0.5),
             size=18, bold=True, color=GOLD, font=FONT_H)
    stack = [
        ("数据库",     "MySQL 8.0  ·  InnoDB 引擎",        COFFEE),
        ("ORM",       "SQLAlchemy 2.x  ·  Flask-SQLAlchemy", GOLD),
        ("后端框架",   "Flask 3.0  ·  应用工厂 + Blueprint", COFFEE_LIGHT),
        ("模板引擎",   "Jinja2  ·  服务端渲染",             BLUE),
        ("前端",      "原生 HTML + CSS  ·  少量 JS",         GREEN),
        ("驱动",      "PyMySQL  +  cryptography",          RED),
        ("版本控制",   "Git + GitHub  ·  Conventional Commits", COFFEE),
    ]
    for i, (k, v, c) in enumerate(stack):
        top = Inches(2.7) + Inches(0.55) * i
        add_round(s, Inches(0.85), top + Inches(0.08),
                  Inches(0.12), Inches(0.32), c, radius=0.4)
        add_text(s, k, Inches(1.1), top, Inches(1.6), Inches(0.45),
                 size=13, bold=True, color=COFFEE_DARK, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, v, Inches(2.75), top, Inches(3.6), Inches(0.45),
                 size=12, color=TEXT_SOFT, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 右：架构图（用色块绘制三层）
    card(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8), WHITE)
    add_text(s, "🏛️ 系统架构", Inches(7.1), Inches(2.15),
             Inches(5), Inches(0.5),
             size=18, bold=True, color=GOLD, font=FONT_H)

    # 三层
    def layer(top, label, sub, color):
        add_round(s, Inches(7.1), top, Inches(5.4), Inches(0.9),
                  color, radius=0.1)
        add_text(s, label, Inches(7.3), top + Inches(0.05),
                 Inches(5), Inches(0.45),
                 size=15, bold=True, color=CREAM, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sub, Inches(7.3), top + Inches(0.45),
                 Inches(5), Inches(0.4),
                 size=11, color=CREAM, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)

    layer(Inches(2.85), "🌐  浏览器（表示层）",
          "Jinja2 模板 · HTML/CSS/JS · 响应式",        COFFEE)
    # 箭头
    add_text(s, "▼", Inches(7.1), Inches(3.78),
             Inches(5.4), Inches(0.3),
             size=18, color=GOLD, align=PP_ALIGN.CENTER, font=FONT_H)
    add_text(s, "HTTP / 表单提交",
             Inches(7.1), Inches(3.92), Inches(5.4), Inches(0.25),
             size=10, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_B)

    layer(Inches(4.2), "⚙️  Flask 应用（业务逻辑层）",
          "11 个 Blueprint · 装饰器鉴权 · SQLAlchemy ORM", COFFEE_LIGHT)
    add_text(s, "▼", Inches(7.1), Inches(5.13),
             Inches(5.4), Inches(0.3),
             size=18, color=GOLD, align=PP_ALIGN.CENTER, font=FONT_H)
    add_text(s, "SQL (经 PyMySQL 驱动)",
             Inches(7.1), Inches(5.27), Inches(5.4), Inches(0.25),
             size=10, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_B)

    layer(Inches(5.55), "🗄️  MySQL（数据持久层）",
          "10 张表 · 3 触发器 · 1 存储过程 · 2 视图",   GOLD)

    footer(s, 5, TOTAL)


# ====================================================================
# 新增：表关系图谱 + 三类关系详解
# ====================================================================
def slide_relations():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·补1", "表关系图谱  —  9 条外键串起 10 张表",
              "ER 模型在数据库中如何具体落地为外键关系")

    # 上半：FK 关系清单
    card(s, Inches(0.5), Inches(2.0), Inches(7.3), Inches(4.9), WHITE)
    add_text(s, "🔗 全部外键关系", Inches(0.75), Inches(2.1),
             Inches(5), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)
    fks = [
        ("products.category_id",         "→",  "categories.category_id",  "商品归属分类"),
        ("orders.customer_id",            "→",  "customers.customer_id",   "订单归属会员（可空）"),
        ("orders.staff_id",               "→",  "staff.staff_id",          "订单创建员工"),
        ("orders.promo_id",               "→",  "promotions.promo_id",     "订单使用的促销码"),
        ("order_items.order_id",          "→",  "orders.order_id",         "订单项归属订单"),
        ("order_items.product_id",        "→",  "products.product_id",     "订单项指向商品"),
        ("product_ingredient.product_id", "→",  "products.product_id",     "配方-商品（联合 PK）"),
        ("product_ingredient.ingredient_id","→","inventory.ingredient_id", "配方-原料（联合 PK）"),
        ("customers.member_level",        "→",  "member_discount.level",   "会员等级（逻辑外键）"),
    ]
    for i, (src, arrow, dst, desc) in enumerate(fks):
        top = Inches(2.55) + Inches(0.43) * i
        # 序号点
        add_round(s, Inches(0.85), top + Inches(0.12),
                  Inches(0.18), Inches(0.18),
                  COFFEE if i % 2 == 0 else GOLD, radius=0.5)
        add_text(s, src, Inches(1.1), top,
                 Inches(2.6), Inches(0.4),
                 size=10, color=COFFEE_DARK, font=FONT_M,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, arrow, Inches(3.7), top, Inches(0.3), Inches(0.4),
                 size=12, bold=True, color=GOLD,
                 font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, dst, Inches(4.0), top,
                 Inches(2.4), Inches(0.4),
                 size=10, color=COFFEE_DARK, font=FONT_M,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(6.4), top,
                 Inches(1.4), Inches(0.4),
                 size=10, color=GRAY, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 右侧：三类关系实例
    card(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.9), WHITE)
    add_text(s, "📐 三类关系", Inches(8.25), Inches(2.1),
             Inches(4), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)

    rel_blocks = [
        ("1 : 1", "  一对一",
         "customers ↔ member_discount\n（每位会员对应唯一等级）",
         RGBColor(0xD3, 0xE3, 0xF4), BLUE),
        ("1 : N", "  一对多",
         "categories → products\norders → order_items\nstaff → orders",
         RGBColor(0xD8, 0xEC, 0xD1), GREEN),
        ("M : N", "  多对多",
         "products ↔ inventory\n通过联系表 product_ingredient\n（每个商品消耗多种原料）",
         RGBColor(0xFC, 0xEF, 0xCB), RGBColor(0x8A, 0x67, 0x00)),
    ]
    for i, (k, name, body, bg, fg) in enumerate(rel_blocks):
        top = Inches(2.65) + Inches(1.4) * i
        add_round(s, Inches(8.25), top, Inches(4.3), Inches(1.25),
                  bg, radius=0.08)
        add_text(s, k, Inches(8.4), top + Inches(0.05),
                 Inches(1.0), Inches(0.35),
                 size=14, bold=True, color=fg, font=FONT_M,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, Inches(9.4), top + Inches(0.05),
                 Inches(2.5), Inches(0.35),
                 size=12, bold=True, color=fg, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, Inches(8.4), top + Inches(0.42),
                 Inches(4.05), Inches(0.82),
                 size=10, color=COFFEE_DARK, font=FONT_B,
                 line_spacing=1.3)
    footer(s, 9, TOTAL)


# ====================================================================
# 新增：完整性约束 + 范式分析
# ====================================================================
def slide_integrity():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·补2", "完整性约束体系与范式分析",
              "四类完整性 · 1NF → BCNF 演化路径")

    # 左：四类完整性
    card(s, Inches(0.5), Inches(2.0), Inches(6.05), Inches(4.9), WHITE)
    add_text(s, "🛡️ 四类完整性约束",
             Inches(0.75), Inches(2.1), Inches(5), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)

    integrities = [
        ("实体完整性",     "Entity",      "所有表都有 PRIMARY KEY，确保元组唯一",
         "PK: customer_id, order_id, ...",                                       COFFEE),
        ("参照完整性",     "Referential", "9 条 FOREIGN KEY 约束，避免悬空引用",
         "ON DELETE SET NULL · CASCADE",                                          GOLD),
        ("域完整性",       "Domain",      "ENUM / CHECK / NOT NULL / DEFAULT 限定取值",
         "status ∈ ENUM('PENDING',...,'CANCELLED')",                              COFFEE_LIGHT),
        ("用户自定义完整性","User-Defined","触发器 · 存储过程实现业务规则",
         "PAID 自动入账积分 · 自动扣库存",                                          BLUE),
    ]
    for i, (zh, en, desc, sample, c) in enumerate(integrities):
        top = Inches(2.65) + Inches(1.0) * i
        add_round(s, Inches(0.75), top + Inches(0.05),
                  Inches(0.12), Inches(0.85), c, radius=0.4)
        add_text(s, zh, Inches(1.0), top,
                 Inches(2.0), Inches(0.4),
                 size=13, bold=True, color=COFFEE_DARK, font=FONT_H)
        add_text(s, en, Inches(3.0), top,
                 Inches(3.0), Inches(0.4),
                 size=10, color=c, font=FONT_M,
                 anchor=MSO_ANCHOR.BOTTOM)
        add_text(s, desc, Inches(1.0), top + Inches(0.38),
                 Inches(5.0), Inches(0.35),
                 size=11, color=TEXT_SOFT, font=FONT_B)
        add_text(s, "  ► " + sample, Inches(1.0), top + Inches(0.7),
                 Inches(5.0), Inches(0.3),
                 size=10, color=c, font=FONT_M, bold=True)

    # 右：范式分析
    card(s, Inches(6.75), Inches(2.0), Inches(6.05), Inches(4.9), WHITE)
    add_text(s, "📚 范式分析  (以 order_items 为例)",
             Inches(7.0), Inches(2.1), Inches(6), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)

    nfs = [
        ("1NF", "第一范式",
         "属性原子化，每个字段不可再分。\n所有表均满足。",
         RGBColor(0xFC, 0xEF, 0xCB)),
        ("2NF", "第二范式",
         "在 1NF 基础上，非主属性完全依赖主键。\n"
         "order_items(item_id PK)，quantity/unit_price 完全依赖主键。",
         RGBColor(0xD8, 0xEC, 0xD1)),
        ("3NF", "第三范式",
         "在 2NF 基础上，消除传递依赖。\n"
         "unit_price 直接随订单项存储而非依赖 product_id "
         "→ 因为商品涨价不应改变历史订单价格（业务需要）。",
         RGBColor(0xD3, 0xE3, 0xF4)),
        ("BCNF","Boyce-Codd",
         "所有决定因素都是候选键。\n"
         "联系表 product_ingredient 的复合主键 "
         "(product_id, ingredient_id) 满足 BCNF。",
         RGBColor(0xF3, 0xD7, 0xD2)),
    ]
    for i, (k, name, desc, bg) in enumerate(nfs):
        top = Inches(2.65) + Inches(1.0) * i
        add_round(s, Inches(7.0), top, Inches(5.65), Inches(0.9),
                  bg, radius=0.08)
        add_text(s, k, Inches(7.15), top + Inches(0.1),
                 Inches(0.75), Inches(0.3),
                 size=14, bold=True, color=COFFEE_DARK,
                 font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, Inches(7.15), top + Inches(0.42),
                 Inches(0.9), Inches(0.3),
                 size=10, color=TEXT_SOFT, font=FONT_B)
        add_text(s, desc, Inches(8.1), top + Inches(0.05),
                 Inches(4.5), Inches(0.85),
                 size=10, color=COFFEE_DARK, font=FONT_B,
                 line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 10, TOTAL)


# ====================================================================
# 新增：索引设计与查询优化
# ====================================================================
def slide_index_design():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·补3", "索引设计与查询优化",
              "主键索引 · 唯一索引 · 外键索引 · 业务热点索引")

    # 上：索引清单表格
    card(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.0), WHITE)
    add_text(s, "🔎 索引清单",
             Inches(0.8), Inches(2.1), Inches(5), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)
    rows = [
        ("类型",        "表 / 列",                              "用途",                          "颜色"),
        ("PRIMARY",     "10 张表的 *_id 列",                    "唯一标识 + 聚簇索引",            COFFEE),
        ("UNIQUE",      "customers.phone · staff.username · promotions.code", "防重复 + 加速查找", GOLD),
        ("FK 索引",     "MySQL 自动为外键创建",                  "支持 JOIN 查询",                COFFEE_LIGHT),
        ("idx_time",    "orders.order_time",                    "仪表盘按日期范围筛选",          BLUE),
        ("idx_status",  "orders.status",                        "订单列表按状态筛选",            GREEN),
        ("复合 PK",     "product_ingredient(product_id, ingredient_id)", "M:N 联系表双向查询", RED),
    ]
    table = s.shapes.add_table(rows=len(rows), cols=3,
                               left=Inches(0.8), top=Inches(2.55),
                               width=Inches(11.7),
                               height=Inches(2.35)).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(5.2)
    table.columns[2].width = Inches(4.7)
    for i, row in enumerate(rows):
        for j in range(3):
            cell = table.cell(i, j); cell.text = row[j]
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT_H if i == 0 else (
                        FONT_M if j != 2 else FONT_B)
                    r.font.size = Pt(12)
                    r.font.bold = (i == 0)
                    r.font.color.rgb = CREAM if i == 0 else COFFEE_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (COFFEE if i == 0 else
                (RGBColor(0xFA, 0xF4, 0xEC) if i % 2 else WHITE))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 下：典型查询 + 用到的索引
    card(s, Inches(0.5), Inches(5.05), Inches(7.3), Inches(1.95), WHITE)
    add_text(s, "📈 典型查询 → 使用的索引",
             Inches(0.75), Inches(5.15), Inches(5), Inches(0.4),
             size=13, bold=True, color=GOLD, font=FONT_H)
    examples = [
        ("当日营收",       "WHERE DATE(order_time)=CURDATE()",  "idx_time",        COFFEE),
        ("待处理订单",     "WHERE status='PENDING'",            "idx_status",      GREEN),
        ("会员手机号查找", "WHERE phone='138...'",              "UNIQUE(phone)",   GOLD),
        ("订单详情 JOIN",  "ON oi.order_id=o.order_id",         "FK idx + PK",     BLUE),
    ]
    for i, (name, cond, idx, c) in enumerate(examples):
        top = Inches(5.55) + Inches(0.35) * i
        add_text(s, "● " + name, Inches(0.8), top,
                 Inches(1.6), Inches(0.3),
                 size=11, bold=True, color=c, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cond, Inches(2.5), top,
                 Inches(3.0), Inches(0.3),
                 size=10, color=COFFEE_DARK, font=FONT_M,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, "→ " + idx, Inches(5.5), top,
                 Inches(2.0), Inches(0.3),
                 size=10, color=c, font=FONT_M, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 右下：性能对比小卡
    card(s, Inches(7.95), Inches(5.05), Inches(4.85), Inches(1.95), WHITE)
    add_text(s, "⚡ 索引带来的提速（估算）",
             Inches(8.2), Inches(5.15), Inches(5), Inches(0.4),
             size=13, bold=True, color=GOLD, font=FONT_H)
    perfs = [
        ("无索引",      "全表扫描 O(N)",   "~150 行扫描",     RED),
        ("PK 主键索引", "B+ 树查找 O(log N)", "~7 次定位",     GREEN),
        ("范围 + 索引", "索引区间扫描",    "毫秒级响应",       BLUE),
    ]
    for i, (k, how, est, c) in enumerate(perfs):
        top = Inches(5.55) + Inches(0.42) * i
        add_round(s, Inches(8.2), top + Inches(0.08),
                  Inches(0.12), Inches(0.22), c, radius=0.4)
        add_text(s, k, Inches(8.4), top,
                 Inches(1.4), Inches(0.35),
                 size=11, bold=True, color=COFFEE_DARK, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, how, Inches(9.85), top,
                 Inches(1.9), Inches(0.35),
                 size=10, color=TEXT_SOFT, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, est, Inches(11.75), top,
                 Inches(1.05), Inches(0.35),
                 size=10, color=c, font=FONT_M, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 11, TOTAL)


# ====================================================================
# 新增：订单状态机 + 折扣业务
# ====================================================================
def slide_business_logic():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "03·6", "核心业务逻辑",
              "订单状态机 · 折扣 ⊕ 积分双重抵扣计算优先级")

    # 上：订单状态机
    card(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.2), WHITE)
    add_text(s, "🔄 订单状态机", Inches(0.8), Inches(2.15),
             Inches(5), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)

    states = [
        ("PENDING",   "待付款",  RGBColor(0xFC, 0xEF, 0xCB), RGBColor(0x8A, 0x67, 0x00)),
        ("PAID",      "已付款",  RGBColor(0xD8, 0xEC, 0xD1), RGBColor(0x2B, 0x6A, 0x2B)),
        ("PREPARING", "制作中",  RGBColor(0xD3, 0xE3, 0xF4), RGBColor(0x25, 0x57, 0x8A)),
        ("READY",     "待取餐",  RGBColor(0xD3, 0xE3, 0xF4), RGBColor(0x25, 0x57, 0x8A)),
        ("COMPLETED", "已完成",  RGBColor(0xD8, 0xEC, 0xD1), RGBColor(0x2B, 0x6A, 0x2B)),
    ]
    bw, bh = Inches(1.85), Inches(0.85)
    base_left = Inches(0.85)
    for i, (en, zh, bg, fg) in enumerate(states):
        left = base_left + (bw + Inches(0.45)) * i
        add_round(s, left, Inches(2.85), bw, bh, bg, radius=0.18)
        add_text(s, en, left, Inches(2.92),
                 bw, Inches(0.35),
                 size=12, bold=True, color=fg,
                 align=PP_ALIGN.CENTER, font=FONT_M)
        add_text(s, zh, left, Inches(3.27),
                 bw, Inches(0.35),
                 size=11, color=fg,
                 align=PP_ALIGN.CENTER, font=FONT_H)
        # 箭头
        if i < len(states) - 1:
            arrow_left = left + bw + Emu(38100)
            add_text(s, "➜", arrow_left, Inches(2.85),
                     Inches(0.42), bh,
                     size=20, color=COFFEE, align=PP_ALIGN.CENTER,
                     font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)

    # CANCELLED 单独
    add_round(s, Inches(5.5), Inches(3.85),
              Inches(2.3), Inches(0.4),
              RGBColor(0xF3, 0xD7, 0xD2), radius=0.3)
    add_text(s, "✕  CANCELLED  已取消",
             Inches(5.5), Inches(3.85), Inches(2.3), Inches(0.4),
             size=11, bold=True, color=RGBColor(0xA2, 0x3A, 0x25),
             align=PP_ALIGN.CENTER, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "（任意中间态均可跳转）",
             Inches(4.4), Inches(3.85), Inches(4.5), Inches(0.4),
             size=10, color=GRAY, font=FONT_B,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # 触发器位置
    add_text(s, "★ PENDING → PAID 触发器自动入账积分",
             Inches(0.8), Inches(3.85), Inches(4.5), Inches(0.4),
             size=10, color=COFFEE, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)

    # 下：折扣计算流程
    card(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6), WHITE)
    add_text(s, "💰 折扣 + 积分 计算优先级（结账时）",
             Inches(0.8), Inches(4.55), Inches(8), Inches(0.4),
             size=15, bold=True, color=GOLD, font=FONT_H)

    # 计算步骤
    steps = [
        ("①", "原价",       "subtotal = Σ(quantity × unit_price)",       COFFEE),
        ("②", "取较优折扣",  "rate = MIN(member_discount, promo_discount)", GOLD),
        ("③", "折后金额",   "discounted = subtotal × rate",              COFFEE_LIGHT),
        ("④", "积分抵现",   "total = MAX(0, discounted − points_used × 0.01)", BLUE),
    ]
    for i, (no, name, formula, c) in enumerate(steps):
        top = Inches(5.05) + Inches(0.45) * i
        add_round(s, Inches(0.85), top, Inches(0.4), Inches(0.4),
                  c, radius=0.25)
        add_text(s, no, Inches(0.85), top, Inches(0.4), Inches(0.4),
                 size=14, bold=True, color=CREAM,
                 align=PP_ALIGN.CENTER, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, Inches(1.35), top,
                 Inches(1.5), Inches(0.4),
                 size=13, bold=True, color=COFFEE_DARK, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, formula, Inches(2.85), top,
                 Inches(9.8), Inches(0.4),
                 size=12, color=TEXT_SOFT, font=FONT_M,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 14, TOTAL)


# ====================================================================
# 新增：角色权限对比
# ====================================================================
def slide_roles():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "04·补", "角色权限矩阵（RBAC）",
              "三种角色 · 装饰器 @admin_required 控制访问")

    roles = [
        ("ADMIN",   "管理员", RED, "@admin_required"),
        ("BARISTA", "咖啡师", GREEN, "@login_required"),
        ("CASHIER", "收银员", BLUE, "@login_required"),
    ]

    # 三栏角色卡片
    cw = Inches(4.05)
    perms = {
        "查看仪表盘 / 统计":    [True, True, True],
        "下单 / 修改订单状态":  [True, True, True],
        "管理菜单（增删改）":   [True, False, False],
        "管理库存 / 配方":      [True, False, False],
        "管理促销 / 会员等级":  [True, False, False],
        "管理员工账号":         [True, False, False],
        "导出 / 系统设置":     [True, False, False],
    }

    for i, (en, zh, color, decor) in enumerate(roles):
        left = Inches(0.5) + (cw + Inches(0.18)) * i
        card(s, left, Inches(2.0), cw, Inches(4.9), WHITE)
        # 顶部色条
        add_rect(s, left, Inches(2.0), cw, Inches(0.7), fill=color)
        add_text(s, zh, left, Inches(2.05), cw, Inches(0.4),
                 size=22, bold=True, color=CREAM,
                 align=PP_ALIGN.CENTER, font=FONT_H)
        add_text(s, en, left, Inches(2.4), cw, Inches(0.3),
                 size=10, color=CREAM,
                 align=PP_ALIGN.CENTER, font=FONT_M)
        # 装饰器名
        add_round(s, left + Inches(0.5), Inches(2.85),
                  cw - Inches(1.0), Inches(0.42),
                  RGBColor(0xFA, 0xF4, 0xEC), radius=0.3)
        add_text(s, decor, left + Inches(0.5), Inches(2.85),
                  cw - Inches(1.0), Inches(0.42),
                  size=11, color=color, font=FONT_M, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 权限列表
        for j, (perm, allowed) in enumerate(perms.items()):
            top = Inches(3.5) + Inches(0.42) * j
            icon = "✓" if allowed[i] else "✕"
            ic = GREEN if allowed[i] else RED
            add_text(s, icon, left + Inches(0.4), top,
                     Inches(0.4), Inches(0.35),
                     size=16, bold=True, color=ic, font=FONT_H,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, perm, left + Inches(0.9), top,
                     cw - Inches(1.1), Inches(0.35),
                     size=11.5,
                     color=(COFFEE_DARK if allowed[i] else GRAY),
                     font=FONT_B, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 16, TOTAL)


# ====================================================================
# 新增：性能优化前后
# ====================================================================
def slide_perf():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "06·补", "性能优化  —  N+1 查询消除",
              "joinedload / selectinload  ·  单页 SQL 数  ~70 → ~5")

    # 左：优化前
    card(s, Inches(0.5), Inches(2.0), Inches(6.05), Inches(4.9), WHITE)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.05), Inches(0.6),
             fill=RED)
    add_text(s, "❌  优化前", Inches(0.8), Inches(2.0),
             Inches(5.5), Inches(0.6),
             size=18, bold=True, color=CREAM, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    code_bad = (
        "# orders 列表：30 条订单 / 页\n"
        "orders = Order.query.all()\n"
        "for o in orders:\n"
        "    print(o.customer.name)   # → 30 次查询\n"
        "    print(o.staff.name)      # → 30 次查询\n"
        "    for it in o.items:       # → 30 次查询\n"
        "        ...\n"
        "# 合计 ≈  1 + 30 × 3  =  91 条 SQL"
    )
    box = s.shapes.add_textbox(Inches(0.75), Inches(2.85),
                               Inches(5.6), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_bad.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run(); r.text = line
        r.font.name = FONT_M; r.font.size = Pt(11.5)
        r.font.color.rgb = (RED if line.startswith("#") else COFFEE_DARK)

    add_round(s, Inches(0.75), Inches(5.5), Inches(5.6), Inches(1.2),
              RGBColor(0xFB, 0xE5, 0xE1), radius=0.08)
    add_text(s, "实测", Inches(0.95), Inches(5.6),
             Inches(2), Inches(0.4), size=12, color=RED, bold=True,
             font=FONT_H)
    add_text(s, "≈  70 条 SQL / 次",
             Inches(0.95), Inches(5.9), Inches(5), Inches(0.45),
             size=22, bold=True, color=RED, font=FONT_H)
    add_text(s, "页面渲染慢，DB 压力大",
             Inches(0.95), Inches(6.35), Inches(5), Inches(0.3),
             size=11, color=GRAY, font=FONT_B)

    # 右：优化后
    card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(4.9), WHITE)
    add_rect(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(0.6),
             fill=GREEN)
    add_text(s, "✓  优化后", Inches(7.0), Inches(2.0),
             Inches(5.5), Inches(0.6),
             size=18, bold=True, color=CREAM, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    code_good = (
        "# 用 joinedload / selectinload 预加载\n"
        "orders = (Order.query\n"
        "    .options(\n"
        "        joinedload(Order.customer),\n"
        "        joinedload(Order.staff),\n"
        "        selectinload(Order.items))\n"
        "    .all())\n"
        "# 合计 =  1 + 1 + 1 + 1  ≈  5 条 SQL"
    )
    box = s.shapes.add_textbox(Inches(6.95), Inches(2.85),
                               Inches(5.7), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_good.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run(); r.text = line
        r.font.name = FONT_M; r.font.size = Pt(11.5)
        r.font.color.rgb = (GREEN if line.startswith("#") else COFFEE_DARK)

    add_round(s, Inches(6.95), Inches(5.5), Inches(5.7), Inches(1.2),
              RGBColor(0xD8, 0xEC, 0xD1), radius=0.08)
    add_text(s, "实测", Inches(7.15), Inches(5.6),
             Inches(2), Inches(0.4), size=12, color=GREEN, bold=True,
             font=FONT_H)
    add_text(s, "≈  5 条 SQL / 次   ⚡ 提速 14×",
             Inches(7.15), Inches(5.9), Inches(5), Inches(0.45),
             size=22, bold=True, color=GREEN, font=FONT_H)
    add_text(s, "页面瞬开，DB 压力大幅下降",
             Inches(7.15), Inches(6.35), Inches(5), Inches(0.3),
             size=11, color=GRAY, font=FONT_B)
    footer(s, 22, TOTAL)


# ====================================================================
# 新增：数据规模 + 学习收获
# ====================================================================
def slide_stats_and_learning():
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, fill=CREAM)
    title_bar(s, "06·补", "项目数据与收获",
              "一组直观的数字 · 一份真切的成长")

    # 上半：数字卡片
    nums = [
        ("10",  "张关系表",         COFFEE),
        ("3",   "个触发器",         GOLD),
        ("29",  "款商品（AI 配图）", COFFEE_LIGHT),
        ("20",  "位演示会员",       BLUE),
        ("150", "条历史订单",       GREEN),
        ("11",  "个 Blueprint",     RED),
        ("8",   "个多表 JOIN 报表", COFFEE),
        ("14×", "性能提升",         GOLD),
    ]
    cw = Inches(1.49)
    for i, (n, label, c) in enumerate(nums):
        left = Inches(0.5) + (cw + Inches(0.05)) * i
        card(s, left, Inches(2.0), cw, Inches(1.7), WHITE)
        add_rect(s, left, Inches(2.0), cw, Inches(0.1), fill=c)
        add_text(s, n, left, Inches(2.2), cw, Inches(0.85),
                 size=32, bold=True, color=c,
                 align=PP_ALIGN.CENTER, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label, left, Inches(3.05), cw, Inches(0.5),
                 size=11, color=TEXT_SOFT,
                 align=PP_ALIGN.CENTER, font=FONT_B)

    # 下半：学习收获
    card(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(2.95), WHITE)
    add_text(s, "🎓 学习收获", Inches(0.8), Inches(4.05),
             Inches(5), Inches(0.4),
             size=18, bold=True, color=GOLD, font=FONT_H)

    learnings = [
        ("数据库设计",  "从需求到 ER 到 3NF 关系的完整建模流程；触发器、视图、存储过程在业务规则中的恰当落地。", COFFEE),
        ("ORM 与 SQL", "理解 ORM 屏蔽下的实际 SQL；通过 EXPLAIN 与日志识别 N+1，并用预加载消除。",          GOLD),
        ("工程实践",   "Blueprint 模块化、应用工厂、装饰器鉴权；数据库改动配套幂等迁移脚本。",              COFFEE_LIGHT),
        ("协作能力",   "Git 分支管理、Conventional Commits、Code Review；二人远程协作的同步与冲突处理。",     BLUE),
    ]
    for i, (k, v, c) in enumerate(learnings):
        top = Inches(4.6) + Inches(0.55) * i
        add_round(s, Inches(0.8), top + Inches(0.1),
                  Inches(0.12), Inches(0.32), c, radius=0.4)
        add_text(s, k, Inches(1.05), top,
                 Inches(2.0), Inches(0.5),
                 size=13, bold=True, color=COFFEE_DARK, font=FONT_H,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, v, Inches(3.1), top,
                 Inches(9.6), Inches(0.5),
                 size=11.5, color=TEXT_SOFT, font=FONT_B,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 23, TOTAL)


# ============================================================
# 构建
# ============================================================
TOTAL = 26

slide_cover()                       #  1
slide_outline()                     #  2
slide_division()                    #  3
slide_overview()                    #  4
slide_techstack()                   #  5
slide_er()                          #  6
slide_logical()                     #  7
slide_ddl()                         #  8
slide_relations()                   #  9  ← 新增：表关系图谱
slide_integrity()                   # 10  ← 新增：完整性 + 范式
slide_index_design()                # 11  ← 新增：索引设计
slide_triggers()                    # 12
slide_queries()                     # 13
slide_business_logic()              # 14
slide_usage()                       # 15
slide_roles()                       # 16
for i, (fname, name, desc) in enumerate(SHOTS_META):
    slide_screenshot(i + 1, fname, name, desc, page=17 + i)   # 17-21
slide_perf()                        # 22
slide_stats_and_learning()          # 23
slide_engineering()                 # 24
slide_future()                      # 25
slide_thanks()                      # 26

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("WROTE:", OUT)
print("Total slides:", len(prs.slides))
